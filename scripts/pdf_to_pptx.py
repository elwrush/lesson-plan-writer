import sys
from io import BytesIO
from pathlib import Path

import fitz
from pptx import Presentation
from pptx.util import Emu

PDF_PATH = Path(sys.argv[1])
OUTPUT_PATH = Path(sys.argv[2])
TARGET_DPI = 96
TARGET_LONG = 1920

doc = fitz.open(str(PDF_PATH))
page = doc[0]
rect = page.rect

page_w = rect.width
page_h = rect.height

target_w = TARGET_LONG
target_h = int(target_w * page_h / page_w)
zoom = target_w / page_w

mat = fitz.Matrix(zoom, zoom)

SLIDE_W = Emu(int(target_w * 914400 / TARGET_DPI))
SLIDE_H = Emu(int(target_h * 914400 / TARGET_DPI))

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank_layout = prs.slide_layouts[6]

for i in range(page_count := len(doc)):
    pix = doc[i].get_pixmap(matrix=mat)
    img_bytes = pix.tobytes("png")
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.add_picture(BytesIO(img_bytes), 0, 0, SLIDE_W, SLIDE_H)

doc.close()
prs.save(str(OUTPUT_PATH))
print(f"PPTX saved: {OUTPUT_PATH} ({page_count} slides, {target_w}x{target_h})")
